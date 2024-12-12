#!/usr/bin/env python3

"""

    Title :      images_to_slide
    Notes :      Automated routine to stack images to in a given layout to make ppt slide
    Output :     .pptx file
    Author :     Ayan Acharyya
    Started :    Jul 2022
    Examples :   run images_to_slide.py
                Before running this script, run:
                run compute_MZgrad.py --system ayan_pleiades --halo 2392 --upto_kpc 10 --xcol rad --output DD1936,DD1946,DD1956 --weight mass --keep
                run projection_plot.py --system ayan_pleiades --halo 2392 --output DD1936,DD1946,DD1956 --do gas,metal,temp --proj x,y,z --galrad 10 --iscolorlog
                run projection_plot.py --system ayan_pleiades --halo 2392 --output DD1936,DD1946,DD1956 --do vrad --proj x,y,z --galrad 10
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import datetime, time
start_time = time.time()

# -----main code-----------------
if __name__ == '__main__':
    # -----------------change paths here--------------------------
    '''
    output_dir = '/Users/acharyya/Documents/presentations/foggie_group_updates/'
    halo_arr = ['8508', '4123'] # ['2392']
    output_arr = ['RD0042'] # ['DD1936', 'DD1946', 'DD1956']
    quantity_arr = ['gas', 'metal', 'temp', 'vrad', 'Zgrad']
    proj_arr = ['x', 'y', 'z']
    galrad = 300 # 10 or 20
    #output_file = output_dir + halo_arr[0] + '_Zgrad_upto%.1Fkpc_investigation.pptx' % (galrad)
    output_file = output_dir + '8508,4123_Zgrad_upto%.1Fkpc_investigation.pptx' % (galrad)
    noaddedtext = False
    '''
    output_dir = '/Users/acharyya/Documents/writings/papers/FOGGIE_Zgrad/Figs/'
    halo_arr = ['8508']#, '5036', '5016', '4123', '2878', '2392']
    output_arr = ['RD0020', 'RD0027', 'RD0042']
    quantity_arr = ['projection', 'profile', 'distribution']
    proj_arr = ['y']
    galrad = 10
    output_file = output_dir + '3_by_3_snapshots_all_automatic.pptx'
    output_to_galrad_dict = {'RD0042': '28.78', 'RD0027': '14.39', 'RD0020': '9.59', 'RD0031': '17.99'}
    noaddedtext = True

    row_arr = quantity_arr
    col_arr = output_arr
    page_arr = halo_arr

    # ---------------------change margins/sizes here---------------------
    '''
    left_margin_in, top_margin_in = 0.5, 1
    height_inch, width_inch = 3, 3.5
    ppt_height_in, ppt_width_in = 10, 18
    fontsize = 20
    '''
    left_margin_in, top_margin_in, horiz_space_in = 1.2, 0.3, 1
    dpi, ppt_to_physical_ratio = 100, 320
    ppt_height_in, ppt_width_in = 2600 / dpi, 2700 / dpi
    height_inch, width_inch = 2600 / ppt_to_physical_ratio, 2500 / ppt_to_physical_ratio
    fontsize = 20

    # -------------------usually no need to change below this------------------------------
    prs = Presentation()
    prs.slide_height = Inches(ppt_height_in)
    prs.slide_width = Inches(ppt_width_in)

    nrow = len(row_arr)
    ncol = len(col_arr)
    nslide = len(page_arr)

    left = Inches(left_margin_in)
    top = Inches(top_margin_in)

    # ----------------looping over all images-------------------------------
    for i, page in enumerate(page_arr):
        blank_slide_layout = prs.slide_layouts[6] # 6 corresponds to empty slide layout
        slide = prs.slides.add_slide(blank_slide_layout)

        #image_path_root = '/Users/acharyya/Work/astro/foggie_outputs/plots_halo_00' + halo_arr[0] + '/nref11c_nref9f/figs/'
        image_path_root = '/Users/acharyya/Work/astro/foggie_outputs/plots_halo_00' + page + '/nref11c_nref9f/figs/'

        # --------slide header----------
        if not noaddedtext:
            txBox = slide.shapes.add_textbox(Inches(0.1), Inches(0.001), width=Inches(1), height=Inches(0.05))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = page
            p.font.size = Pt(fontsize)
            p.font.bold = True

        for j, row in enumerate(row_arr):
            for k, col in enumerate(col_arr):
                if page == '2878' and col == 'RD0042': col = 'RD0031' # because Cyclone has not run till z=0 yet
                if col_arr == quantity_arr:
                    if col == 'Zgrad': col_path = 'datashader_log_metal_vs_rad_upto%.1Fkpc_wtby_mass.png' % (galrad)
                    else: col_path = col + '_box=%.2Fkpc_proj_%s_projection.png' % (2 * galrad, row)
                elif col_arr == output_arr:
                    if row == 'projection': col_path = col + '_metal_box=' + output_to_galrad_dict[col] + 'kpc_proj_' + proj_arr[0] + '_projection_wdencut.png'
                    elif row == 'profile': col_path = col + '_datashader_log_metal_vs_rad_upto%.1Fckpchinv_wtby_mass_wdencut.png' %(galrad)
                    elif row == 'distribution': col_path = col + '_log_metal_distribution_upto%.1Fckpchinv_wtby_mass_wdencut_islog.png' %(galrad)

                # ----column headers----------
                if j == 0 and not noaddedtext:
                    txBox = slide.shapes.add_textbox(Inches(2) + Inches(k * width_inch), Inches(0.01), width=Inches(width_inch), height=Inches(0.1))
                    tf = txBox.text_frame
                    p = tf.add_paragraph()
                    p.text = col
                    p.font.size = Pt(fontsize)

                # ----rotated row headers--------
                if k == 0 and not noaddedtext:
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(1) + Inches(j * height_inch), width=Inches(0.5), height=Inches(height_inch))
                    txBox.rotation = -90
                    tf = txBox.text_frame
                    p = tf.add_paragraph()
                    p.text = row
                    p.font.size = Pt(fontsize)

                #image_path = image_path_root + page + '/' + col_path
                #image_path = image_path_root + output_arr[0] + '/' + col_path
                image_path = image_path_root +  '/' + col_path
                thisleft = left + Inches(k * width_inch)
                thistop = top + Inches(j * height_inch)

                crop_left_in = 1.3 if j == 0 else 1.0
                crop_bottom_in = 1.1 if j == 0 else 0.
                crop_right_in = 1.6 if j == 0 else 0.1
                crop_top_in = 0.1 if j == 0 else 0.1
                hspace = Inches(horiz_space_in) if j == 0 and k!= 0 else 0.

                pic = slide.shapes.add_picture(image_path, thisleft + hspace, thistop, height=Inches(height_inch))#, width=Inches(width_inch))

                pic.crop_left = crop_left_in / pic.width.inches
                pic.crop_bottom = crop_bottom_in / pic.height.inches
                pic.crop_right = crop_right_in / pic.width.inches
                pic.crop_top = crop_top_in / pic.height.inches

                print('Added', image_path, 'at slide', i+1, 'row', j+1, 'column', k+1, 'which is', i * nrow * ncol + j * ncol + k + 1, 'of', nslide * nrow * ncol, 'images..')

    prs.save(output_file)
    print('Saved files', output_file)

    print('Completed in %s' % (datetime.timedelta(minutes=(time.time() - start_time) / 60)))